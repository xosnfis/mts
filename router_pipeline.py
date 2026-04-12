"""
GPTHub — Smart Model Router + Autonomous Active Memory
OpenWebUI Pipeline Filter.

Memory pipeline:
  inlet  → regex fast-extract facts from last user message
         → semantic search (bge-m3) for top-5 relevant memories
         → inject [USER CONTEXT] block into system prompt
  outlet → mts-anya LLM summarizer extracts/merges facts → SQLite

Multimodal routing:
  Audio  → whisper-turbo-local transcription → inject transcript → route to default/code model
  Vision → detect image_url parts → normalize base64 → route to qwen2.5-vl or cotype-pro-vl-32b
  Files  → detect file attachments → bge-m3 RAG over OpenWebUI collections → inject [FILE CONTEXT]
"""

import re
import json
import base64
import sqlite3
import logging
import asyncio
import httpx
import numpy as np
import sys
from pathlib import Path
from typing import Optional, Callable, Awaitable
from pydantic import BaseModel, Field

# Add /app/utils to path so helper modules don't get picked up as pipelines
_UTILS_PATH = "/app/utils"
if _UTILS_PATH not in sys.path:
    sys.path.insert(0, _UTILS_PATH)

# Multi-format document parser
try:
    from file_processor import extract_text_from_any, build_file_context_block
    _FILE_PROCESSOR_AVAILABLE = True
except ImportError:
    _FILE_PROCESSOR_AVAILABLE = False

# File generator
try:
    from file_generator import generate_file, make_download_link, strip_generate_tag
    _FILE_GENERATOR_AVAILABLE = True
except ImportError:
    _FILE_GENERATOR_AVAILABLE = False

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Storage — SQLite (WAL mode, indexed, per-chat isolation)
#
# Schema:
#   memory_facts  — current facts per (user_id, chat_id, key)
#   memory_history — full audit log of every fact change
# ---------------------------------------------------------------------------

DB_PATH = Path("/app/backend/data/smart_router_memory.db")
_db_lock = asyncio.Lock()


def _db_connect() -> sqlite3.Connection:
    """Open a connection with WAL mode and foreign keys enabled."""
    DB_PATH.parent.mkdir(parents=True, exist_ok=True)
    con = sqlite3.connect(str(DB_PATH), check_same_thread=False)
    con.execute("PRAGMA journal_mode=WAL")
    con.execute("PRAGMA foreign_keys=ON")
    con.execute("PRAGMA synchronous=NORMAL")
    con.row_factory = sqlite3.Row
    return con


def _db_init() -> sqlite3.Connection:
    con = _db_connect()
    con.executescript(
        """
        CREATE TABLE IF NOT EXISTS memory_facts (
            id         INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id    TEXT    NOT NULL,
            chat_id    TEXT    NOT NULL DEFAULT '',
            key        TEXT    NOT NULL,
            value      TEXT    NOT NULL,
            embedding  TEXT,
            updated_at INTEGER NOT NULL DEFAULT (strftime('%s','now')),
            UNIQUE(user_id, chat_id, key)
        );

        CREATE INDEX IF NOT EXISTS idx_facts_user_chat
            ON memory_facts(user_id, chat_id);

        CREATE TABLE IF NOT EXISTS memory_history (
            id         INTEGER PRIMARY KEY AUTOINCREMENT,
            user_id    TEXT    NOT NULL,
            chat_id    TEXT    NOT NULL DEFAULT '',
            key        TEXT    NOT NULL,
            value      TEXT    NOT NULL,
            created_at INTEGER NOT NULL DEFAULT (strftime('%s','now'))
        );

        CREATE INDEX IF NOT EXISTS idx_history_user_chat
            ON memory_history(user_id, chat_id);
        """
    )
    con.commit()
    return con


def _parse_uid(user_id: str) -> tuple[str, str]:
    """Split composite 'uid:chat_id' key back into (uid, chat_id)."""
    if ":" in user_id:
        uid, chat_id = user_id.split(":", 1)
        return uid, chat_id
    return user_id, ""


def _load_all_facts(user_id: str) -> dict:
    uid, chat_id = _parse_uid(user_id)
    try:
        con = _db_init()
        rows = con.execute(
            "SELECT key, value FROM memory_facts WHERE user_id=? AND chat_id=?",
            (uid, chat_id),
        ).fetchall()
        con.close()
        return {r["key"]: r["value"] for r in rows}
    except Exception as e:
        log.error("Memory load error: %s", e)
        return {}


def _save_fact(user_id: str, key: str, value: str, embedding: Optional[list] = None) -> None:
    uid, chat_id = _parse_uid(user_id)
    try:
        con = _db_init()
        emb_str = json.dumps(embedding) if embedding else None
        now = int(__import__("time").time())
        con.execute(
            """INSERT INTO memory_facts(user_id, chat_id, key, value, embedding, updated_at)
               VALUES(?,?,?,?,?,?)
               ON CONFLICT(user_id, chat_id, key) DO UPDATE
               SET value=excluded.value,
                   embedding=excluded.embedding,
                   updated_at=excluded.updated_at""",
            (uid, chat_id, key, value, emb_str, now),
        )
        # Audit trail
        con.execute(
            "INSERT INTO memory_history(user_id, chat_id, key, value, created_at) VALUES(?,?,?,?,?)",
            (uid, chat_id, key, value, now),
        )
        con.commit()
        con.close()
    except Exception as e:
        log.error("Memory save error: %s", e)


def _clear_memory(user_id: str) -> None:
    uid, chat_id = _parse_uid(user_id)
    try:
        con = _db_init()
        con.execute(
            "DELETE FROM memory_facts WHERE user_id=? AND chat_id=?", (uid, chat_id)
        )
        con.commit()
        con.close()
    except Exception as e:
        log.error("Memory clear error: %s", e)


def _load_facts_with_embeddings(user_id: str) -> list[dict]:
    uid, chat_id = _parse_uid(user_id)
    try:
        con = _db_init()
        rows = con.execute(
            """SELECT key, value, embedding FROM memory_facts
               WHERE user_id=? AND chat_id=?
               ORDER BY updated_at DESC""",
            (uid, chat_id),
        ).fetchall()
        con.close()
        result = []
        for r in rows:
            emb = json.loads(r["embedding"]) if r["embedding"] else None
            result.append({"key": r["key"], "value": r["value"], "embedding": emb})
        return result
    except Exception as e:
        log.error("Embedding load error: %s", e)
        return []


# ---------------------------------------------------------------------------
# Embeddings + Semantic Search
# ---------------------------------------------------------------------------

async def _get_embedding(base_url: str, api_key: str, text: str) -> Optional[list]:
    """Get bge-m3 embedding for a text string."""
    try:
        async with httpx.AsyncClient(timeout=15) as client:
            resp = await client.post(
                f"{base_url}/embeddings",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={"model": "bge-m3", "input": text},
            )
        resp.raise_for_status()
        return resp.json()["data"][0]["embedding"]
    except Exception as e:
        log.error("Embedding error: %s", e)
        return None


def _cosine_similarity(a: list, b: list) -> float:
    va, vb = np.array(a, dtype=float), np.array(b, dtype=float)
    denom = np.linalg.norm(va) * np.linalg.norm(vb)
    return float(np.dot(va, vb) / denom) if denom > 0 else 0.0


async def _semantic_search(
    base_url: str, api_key: str, user_id: str, query: str, top_k: int = 5
) -> list[dict]:
    """Return top_k most relevant facts for the query using bge-m3."""
    facts = _load_facts_with_embeddings(user_id)
    if not facts:
        return []

    query_emb = await _get_embedding(base_url, api_key, query)
    if query_emb is None:
        # Fallback: return all facts if embedding fails
        return facts[:top_k]

    scored = []
    for fact in facts:
        if fact["embedding"]:
            score = _cosine_similarity(query_emb, fact["embedding"])
        else:
            score = 0.0
        scored.append((score, fact))

    scored.sort(key=lambda x: x[0], reverse=True)
    return [f for _, f in scored[:top_k]]


# ---------------------------------------------------------------------------
# Regex fast-extraction (runs synchronously, no LLM cost)
# ---------------------------------------------------------------------------

_NAME_STOPWORDS = (
    "работаю", "являюсь", "из", "живу", "нахожусь", "сменил", "теперь",
    "не", "был", "буду", "хочу", "люблю", "думаю", "знаю", "понимаю",
    "могу", "должен", "должна", "пришёл", "пришла", "иду", "иногда",
    "всегда", "никогда", "очень", "просто", "тоже", "уже", "ещё",
    "помню", "помнишь", "скажи", "скажу",
)
_STOPWORDS_RE = "|".join(_NAME_STOPWORDS)


def _extract_facts(text: str) -> dict:
    found: dict = {}

    # Name — RU
    m = re.search(
        r"\b(?:меня\s+зовут|зовут\s+меня|зовут)\s+([а-яёa-z]{2,20})\b",
        text, re.IGNORECASE,
    )
    if m and m.group(1).lower() not in ("меня", "тебя", "его", "её", "нас", "вас"):
        found["name"] = m.group(1).strip()

    # Name — EN
    if "name" not in found:
        m = re.search(r"\bmy name is ([a-z]{2,20})\b", text, re.IGNORECASE)
        if m:
            found["name"] = m.group(1).strip()

    # Name — "я X" heuristic
    if "name" not in found:
        m = re.search(
            r"\bя\s+([а-яё]{2,15})\b(?!\s+(?:" + _STOPWORDS_RE + r"))",
            text, re.IGNORECASE,
        )
        if m:
            candidate = m.group(1).lower()
            if not re.search(r"(ю|ть|ешь|ает|ают|ал|ала|али|ить|еть|лась|лся)$", candidate):
                if candidate not in ("тоже", "сам", "сама", "себя", "очень", "меня"):
                    found["name"] = m.group(1).strip()

    # Age
    m = re.search(r"\bмне\s+(\d{1,3})(?:\s+лет)?\b", text, re.IGNORECASE)
    if m:
        found["age"] = m.group(1)
    if "age" not in found:
        m = re.search(r"\bi(?:'m| am)\s+(\d{1,3})(?:\s+years?\s+old)?\b", text, re.IGNORECASE)
        if m:
            found["age"] = m.group(1)

    # Role / job
    m = re.search(r"\bработаю\s+([\w\s]{2,30}?)(?:\s*[,.]|$)", text, re.IGNORECASE)
    if m:
        found["role"] = m.group(1).strip()
    if "role" not in found:
        m = re.search(r"\bя\s+(?:являюсь\s+)?([а-яё]{3,}(?:ом|ем|ой|ей))\b", text, re.IGNORECASE)
        if m:
            found["role"] = m.group(1).strip()
    if "role" not in found:
        m = re.search(r"\bi\s+work\s+as\s+a(?:n)?\s+([\w\s]{2,30}?)(?:\s*[,.]|$)", text, re.IGNORECASE)
        if m:
            found["role"] = m.group(1).strip()

    # Likes / hobbies
    m = re.search(r"\b(?:люблю|обожаю)\s+([\w\s]{2,40}?)(?:\s*[,.]|$)", text, re.IGNORECASE)
    if m:
        found["likes"] = m.group(1).strip()

    # Location — RU
    m = re.search(r"\bя\s+из\s+([\w\s\-]{2,40}?)(?:\s*[,.]|$)", text, re.IGNORECASE)
    if m:
        found["location"] = m.group(1).strip()
    if "location" not in found:
        m = re.search(r"\b(?:живу|нахожусь)\s+в\s+([\w\s\-]{2,40}?)(?:\s*[,.]|$)", text, re.IGNORECASE)
        if m:
            found["location"] = m.group(1).strip()
    # Location — EN
    if "location" not in found:
        m = re.search(r"\bi\s+live\s+in\s+([\w\s\-]{2,40}?)(?:\s*[,.]|$)", text, re.IGNORECASE)
        if m:
            found["location"] = m.group(1).strip()

    # Project name
    m = re.search(r"\b(?:наш\s+проект|проект)\s+(?:называется|зовётся|это)\s+([\w\s\-]{2,40}?)(?:\s*[,.]|$)", text, re.IGNORECASE)
    if m:
        found["project"] = m.group(1).strip()

    # Tech stack
    m = re.search(r"\b(?:использую|пишу\s+на|работаю\s+с)\s+([\w\s\+\#]{2,30}?)(?:\s*[,.]|$)", text, re.IGNORECASE)
    if m:
        found["tech"] = m.group(1).strip()
    if "tech" not in found:
        m = re.search(r"\bi\s+use\s+([\w\s\+\#]{2,30}?)(?:\s*[,.]|$)", text, re.IGNORECASE)
        if m:
            found["tech"] = m.group(1).strip()

    return found


# ---------------------------------------------------------------------------
# LLM-based summarizer (mts-anya) — runs in outlet, non-blocking
# ---------------------------------------------------------------------------

async def _summarize_and_save(
    base_url: str, api_key: str, model: str, user_id: str, messages: list
) -> None:
    """Use mts-anya to extract facts from conversation and persist with embeddings."""
    user_lines = []
    for m in messages[-20:]:
        if m.get("role") == "user":
            c = m.get("content", "")
            if isinstance(c, str):
                user_lines.append(c)
    if not user_lines:
        return

    conversation = "\n".join(user_lines)
    existing = _load_all_facts(user_id)
    existing_str = json.dumps(existing, ensure_ascii=False) if existing else "{}"

    prompt = (
        "Ты — система памяти. Проанализируй сообщения пользователя и извлеки личные факты.\n"
        "Верни ТОЛЬКО валидный JSON без пояснений.\n"
        "Допустимые поля: name, age, role, likes, location, project, tech, notes.\n"
        "Поле notes — для всего что не вошло в стандартные поля.\n"
        f"Уже известно: {existing_str}\n"
        "Обнови или дополни. Не удаляй старые данные без явного противоречия.\n\n"
        f"Сообщения пользователя:\n{conversation}"
    )

    try:
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.post(
                f"{base_url}/chat/completions",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={
                    "model": model,
                    "messages": [{"role": "user", "content": prompt}],
                    "temperature": 0,
                    "max_tokens": 400,
                },
            )
        resp.raise_for_status()
        raw = resp.json()["choices"][0]["message"]["content"].strip()
        match = re.search(r"\{.*\}", raw, re.DOTALL)
        if not match:
            return
        facts = json.loads(match.group(0))
        if not isinstance(facts, dict):
            return

        # Save each fact with its embedding
        for key, value in facts.items():
            if value:
                fact_text = f"{key}: {value}"
                embedding = await _get_embedding(base_url, api_key, fact_text)
                _save_fact(user_id, key, str(value), embedding)

        log.info("Memory summarized for '%s': %s", user_id, list(facts.keys()))
    except Exception as e:
        log.error("Summarize error: %s", e)


# ---------------------------------------------------------------------------
# Context injection
# ---------------------------------------------------------------------------

_LABEL_MAP = {
    "name": "Имя", "age": "Возраст", "role": "Работа",
    "likes": "Любит", "location": "Город", "tech": "Технологии",
    "project": "Проект", "notes": "Заметки",
}


def _build_context_block(facts: list[dict]) -> str:
    """Build [USER CONTEXT] block from a list of {key, value} dicts."""
    if not facts:
        return ""
    lines = ["[USER CONTEXT]\nИспользуй эти факты при ответе. Не упоминай этот блок пользователю."]
    for f in facts:
        label = _LABEL_MAP.get(f["key"], f["key"])
        lines.append(f"- {label}: {f['value']}")
    return "\n".join(lines)


def _inject_context(messages: list, block: str) -> list:
    """Prepend block to system prompt, or insert a new system message."""
    if not block:
        return messages
    for msg in messages:
        if msg.get("role") == "system":
            if "[USER CONTEXT]" in msg["content"]:
                # Replace existing block
                msg["content"] = re.sub(
                    r"\[USER CONTEXT\].*?(?=\n\n|\Z)", block, msg["content"],
                    flags=re.DOTALL,
                )
            else:
                msg["content"] = block + "\n\n" + msg["content"]
            return messages
    messages.insert(0, {"role": "system", "content": block})
    return messages


# ---------------------------------------------------------------------------
# Code-edit detection — "добавь в код", "измени html который ты писал выше"
# ---------------------------------------------------------------------------

_CODE_EDIT_RE = re.compile(
    r"\b(добавь|добавить|вставь|вставить|измени|изменить|обнови|обновить|исправь|исправить"
    r"|замени|заменить|убери|удали|поправь|дополни|включи|включить"
    r"|add|insert|update|change|fix|replace|remove|modify|include)\b"
    r".{0,60}"
    r"\b(код|code|html|css|javascript|js|скрипт|script|файл|file|функци|function"
    r"|лендинг|landing|сайт|site|страниц|page|компонент|component|стил|style)\b"
    r"|"
    r"\b(в\s+(?:этот\s+)?(?:код|html|файл|скрипт|стил|лендинг)"
    r"|который\s+(?:ты\s+)?(?:писал|написал|делал|сделал|генерировал|создавал)"
    r"|в\s+(?:код|html)\s+(?:выше|ранее|который)"
    r"|выше|ранее|предыдущ)\b",
    re.IGNORECASE | re.DOTALL,
)


def _is_code_edit_request(text: str) -> bool:
    return bool(_CODE_EDIT_RE.search(text))


# ---------------------------------------------------------------------------
# Image-gen detection (mirrors image_gen_pipeline._IMAGE_RE)
# Used after audio transcription to avoid stealing image requests
# ---------------------------------------------------------------------------

_IMAGE_GEN_RE = re.compile(
    r"\b(нарисуй|нарисуйте|нарисуем|нарисую|нарисует|нарисовать|нарисуй-ка"
    r"|изобрази|изобразите|изобразим|изображу|изобразить"
    r"|визуализируй|визуализируйте|визуализируем|визуализировать"
    r"|покажи|покажите|покажем|покажу|покажет"
    r"|набросай|набросок|зарисуй"
    r"|draw|paint|illustrate|sketch|render|depict|visualize|show\s+me|generate\s+image)\b"
    r"|"
    r"\b(сгенерируй|сгенерируем|сгенерирую|сгенерирует|сгенерировать|сгенерируйте"
    r"|создай|создадим|создам|создаст|создать|создайте"
    r"|сделай|сделаем|сделаю|сделает|сделать|сделайте"
    r"|придумай|придумаем|придумаю|придумать|придумайте"
    r"|нужно\s+(?:сгенерировать|создать|нарисовать|сделать|придумать)"
    r"|можешь\s+(?:сгенерировать|нарисовать|создать|сделать)"
    r"|хочу\s+(?:картинку|изображение|фото|логотип|рисунок|арт|иллюстрацию)"
    r"|generate|create|make|produce|design|build|get\s+me)\s+.{0,60}"
    r"(картинк|изображени|фотографи|фото|рисунок|арт\b|иллюстраци|логотип|лого\b|баннер|обложк|аватар|иконк|постер|плакат|обои|превью|миниатюр"
    r"|image|picture|photo|illustration|logo|banner|cover|avatar|icon|poster|artwork|graphic|wallpaper|thumbnail|preview)"
    r"|"
    r"\b(нарисую|нарисуй|сгенерируй|создай|сделай|покажи)\s+(?:мне\s+)?(?:для\s+)?[\w\-]+\.(?:jpg|jpeg|png|webp|gif|svg)\b",
    re.IGNORECASE | re.DOTALL,
)


def _is_image_gen_request(text: str) -> bool:
    return bool(_IMAGE_GEN_RE.search(text))


_PPTX_ROUTER_RE = re.compile(
    r"\bпрезентаци[юяи]\b|\bслайды\b|\bpptx\b|\bpowerpoint\b"
    r"|presentation|make\s+slides|create\s+slides",
    re.IGNORECASE,
)


def _is_pptx_request(text: str) -> bool:
    return bool(_PPTX_ROUTER_RE.search(text))


# ---------------------------------------------------------------------------
# Web-search force-trigger — factual / sourced queries
# ---------------------------------------------------------------------------

_FORCE_SEARCH_RE = re.compile(
    r"\b(источник\w*|ссылк\w+|линк\w*|url\b|cite|citation|source\w*"
    r"|данные\s+(?:за|на|по)\s+\d{4}"
    r"|статистик\w+|исследовани\w+|доказательств\w+"
    r"|по\s+данным|согласно|согласно\s+исследовани\w+"
    r"|актуальн\w+\s+данн\w+|свежи\w+\s+данн\w+"
    r"|в\s+\d{4}\s+году|за\s+\d{4}\s+год"
    r"|links?|sources?|references?|citations?|proof|evidence"
    r"|data\s+(?:from|as\s+of|for)\s+\d{4}"
    r"|latest\s+(?:data|stats|statistics|numbers|figures)"
    r"|according\s+to|based\s+on\s+(?:data|research|study))\b",
    re.IGNORECASE,
)


def _needs_web_search(text: str) -> bool:
    return bool(_FORCE_SEARCH_RE.search(text))


# ---------------------------------------------------------------------------
# File generation trigger
# ---------------------------------------------------------------------------

# ---------------------------------------------------------------------------
# File generation / conversion detection
# ---------------------------------------------------------------------------

# Catches: "сделай pdf", "сохрани как word", "в формате xlsx", "make a pdf", etc.
_FILE_GEN_RE = re.compile(
    r"\b(сохрани|сохранить|запиши|записать|выгрузи|выгрузить|скачай|скачать"
    r"|сделай|сделать|создай|создать|напиши|написать|подготовь|подготовить"
    r"|экспортируй|экспортировать|генерируй|сгенерируй)\b"
    r".{0,40}"
    r"\b(pdf|docx|word|xlsx|excel|csv|txt)\b"
    r"|"
    r"\b(в\s+формате|как\s+файл|в\s+виде\s+файла|в\s+файл)\b"
    r".{0,30}"
    r"\b(pdf|docx|word|xlsx|excel|csv|txt)\b"
    r"|"
    r"\b(pdf|docx|word|xlsx|excel|csv|txt)\s+(файл|документ|отчёт|report|file|doc)\b"
    r"|"
    r"\b(save\s+as|export\s+(?:as|to)|download\s+as|make\s+(?:a\s+)?|create\s+(?:a\s+)?|generate\s+(?:a\s+)?)"
    r"(pdf|docx|word|xlsx|excel|csv|txt)\b",
    re.IGNORECASE | re.DOTALL,
)

# Catches: "конвертируй в pdf", "теперь сделай это word", "convert to excel"
_CONVERT_RE = re.compile(
    r"\b(преобразуй|преобразовать|конвертируй|конвертировать|переведи\s+в"
    r"|сделай\s+(?:из\s+этого|это|этот\s+текст)"
    r"|теперь\s+(?:сделай|создай|сохрани)"
    r"|convert\s+(?:this\s+)?(?:to|into)|turn\s+(?:this\s+)?into)\b"
    r".{0,50}"
    r"\b(pdf|docx|word|xlsx|excel|csv|txt|pptx|презентацию?)\b",
    re.IGNORECASE | re.DOTALL,
)

_FORMAT_DETECT_RE = re.compile(
    r"\b(pdf|docx|word|xlsx|excel|csv|txt|text)\b", re.IGNORECASE
)

# Tag the LLM must append — [ACTION:GENERATE_FILE(type=PDF)]
_ACTION_TAG_RE = re.compile(
    r"\[ACTION:GENERATE_FILE\(type=(\w+)\)\]", re.IGNORECASE
)


def _detect_file_format(text: str) -> Optional[str]:
    """Extract requested output format from user text."""
    m = _FORMAT_DETECT_RE.search(text)
    if not m:
        return None
    fmt = m.group(1).lower()
    return {"docx": "docx", "word": "docx", "pdf": "pdf",
            "xlsx": "xlsx", "excel": "xlsx", "csv": "csv",
            "txt": "txt", "text": "txt"}.get(fmt)


# ---------------------------------------------------------------------------
# PPTX generation (inline — runs inside router inlet)
# ---------------------------------------------------------------------------

_SLIDE_COUNT_RE = re.compile(
    r"\b(один|два|три|четыре|пять|шесть|семь|восемь|девять|десять|\d+)\s+слайд\w*\b"
    r"|\bиз\s+(один|два|три|четыре|пять|шесть|семь|восемь|девять|десять|\d+)\s+слайд\w*\b"
    r"|\bтоп[\s\-]?(один|два|три|четыре|пять|шесть|семь|восемь|девять|десять|\d+)\b"
    r"|\b(один|два|три|четыре|пять|шесть|семь|восемь|девять|десять|\d+)\s+(?:полезн|лучш|топ|интересн|основн|главн|важн)\w*\b",
    re.IGNORECASE,
)


def _pptx_slide_count(text: str) -> int:
    m = _SLIDE_COUNT_RE.search(text)
    if not m:
        return 7
    # Get first non-None group
    val = next((g for g in m.groups() if g is not None), None)
    if val is None:
        return 7
    val = val.lower()
    return int(val) if val.isdigit() else _WORD_TO_NUM.get(val, 7)


def _pptx_topic(text: str) -> str:
    m = re.search(r"на\s+тему\s+[«\"]?(.+?)[»\"]?(?:\.|$)", text, re.IGNORECASE)
    if m:
        return m.group(1).strip(" .,»\"")
    m = re.search(r"(?:презентацию?\s+(?:про|о|об|по)|слайды\s+(?:про|о|об|по))\s+(.+?)(?:\.|$)", text, re.IGNORECASE)
    if m:
        return m.group(1).strip(" .,")
    cleaned = re.sub(
        r"\b(сделай|создай|напиши|сгенерируй|подготовь|составь|сгенерирую|make|create|generate|build)"
        r"\s+(?:мне\s+)?(?:a\s+)?(?:презентацию?|слайды|presentation|pptx|powerpoint|slides)\b",
        "", text, flags=re.IGNORECASE,
    )
    cleaned = re.sub(r"\b(\d+|один|два|три|четыре|пять|шесть|семь|восемь|девять|десять)\s+слайд\w*\b", "", cleaned, flags=re.IGNORECASE)
    cleaned = re.sub(r"\s{2,}", " ", cleaned).strip(" ,.")
    return cleaned or text


async def _pptx_generate_slides(base_url: str, api_key: str, model: str, topic: str, num: int, web_context: str = "") -> list[dict]:
    context_block = f"\n\nДополнительный контекст из интернета:\n{web_context}" if web_context else ""
    prompt = (
        f"Создай содержимое для презентации на тему: «{topic}».\n"
        f"Нужно {num} слайдов (не считая титульного).{context_block}\n"
        f"Верни ТОЛЬКО валидный JSON-массив без пояснений. Формат:\n"
        f'[{{"title": "Заголовок", "bullets": ["Пункт 1", "Пункт 2", "Пункт 3"]}}, ...]\n'
        f"Каждый слайд: 1 заголовок + 3-5 коротких пунктов (до 15 слов каждый). Язык: русский."
    )
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(
                f"{base_url}/chat/completions",
                headers={"Authorization": f"Bearer {api_key}"},
                json={"model": model, "messages": [{"role": "user", "content": prompt}],
                      "temperature": 0.4, "max_tokens": 1500},
            )
        r.raise_for_status()
        raw = r.json()["choices"][0]["message"]["content"].strip()
        m = re.search(r"\[.*\]", raw, re.DOTALL)
        if m:
            return json.loads(m.group(0))[:num]
    except Exception as e:
        log.error("PPTX slide gen error: %s", e)
    return [{"title": f"Раздел {i}", "bullets": [f"Содержимое {i}"]} for i in range(1, num + 1)]


def _pptx_build(topic: str, slides_data: list[dict], output_path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    BG = RGBColor(0x1A, 0x1A, 0x2E)
    HIGHLIGHT = RGBColor(0x0F, 0x3C, 0x78)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xCC, 0xDD, 0xFF)
    GOLD = RGBColor(0xFF, 0xD7, 0x00)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def _bg(slide, color):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = color

    def _tb(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
        tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def _rect(slide, l, t, w, h, color):
        s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = color
        s.line.fill.background()

    # Title slide
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG)
    _rect(sl, 0, 2.8, 13.33, 0.06, GOLD)
    _tb(sl, topic, 1.0, 1.2, 11.33, 1.8, size=40, bold=True, align=PP_ALIGN.CENTER)
    _tb(sl, "Выполнил/а:", 1.0, 3.2, 11.33, 0.5, size=22, color=LIGHT, align=PP_ALIGN.CENTER)
    _tb(sl, "Проверил/а", 1.0, 3.75, 11.33, 0.5, size=22, color=LIGHT, align=PP_ALIGN.CENTER)

    # Content slides
    for i, sd in enumerate(slides_data, 2):
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        _bg(sl, BG)
        _rect(sl, 0, 0, 13.33, 1.3, HIGHLIGHT)
        _tb(sl, sd.get("title", ""), 0.4, 0.15, 12.5, 1.0, size=28, bold=True)
        _rect(sl, 0.4, 1.45, 12.5, 0.04, GOLD)
        y = 1.65
        for b in sd.get("bullets", []):
            _tb(sl, f"▸  {b}", 0.6, y, 12.0, 0.55, size=18, color=LIGHT)
            y += 0.65

    prs.save(output_path)


def _pptx_html_preview(topic: str, slides_data: list[dict]) -> str:
    slides_html = [f"""  <div class="slide title-slide">
    <div class="sn">1</div>
    <h1>{topic}</h1>
    <div class="meta"><span>Выполнил/а</span><span>Проверил/а</span></div>
  </div>"""]
    for i, s in enumerate(slides_data, 2):
        bullets = "\n".join(f"    <li>{b}</li>" for b in s.get("bullets", []))
        slides_html.append(f"""  <div class="slide content-slide">
    <div class="sn">{i}</div>
    <h2>{s.get('title','')}</h2>
    <div class="div"></div>
    <ul>{bullets}</ul>
  </div>""")
    return """<!DOCTYPE html><html lang="ru"><head><meta charset="UTF-8"><style>
body{font-family:'Segoe UI',Arial,sans-serif;background:#0d0d1a;margin:0;padding:16px;display:flex;flex-direction:column;gap:16px}
.slide{max-width:860px;border-radius:10px;position:relative;overflow:hidden;box-shadow:0 4px 20px rgba(0,0,0,.5)}
.title-slide{background:linear-gradient(135deg,#1a1a2e,#16213e,#0f3c78);min-height:240px;padding:32px 40px}
.content-slide{background:linear-gradient(160deg,#1a1a2e,#16213e);padding:32px 40px 28px}
.sn{position:absolute;top:12px;right:16px;color:rgba(255,215,0,.5);font-size:12px;font-weight:600}
h1{color:#fff;font-size:28px;font-weight:700;margin:20px 0 24px}
.meta{display:flex;gap:32px;color:#aabbdd;font-size:14px;border-top:2px solid #ffd700;padding-top:14px}
h2{color:#fff;font-size:20px;font-weight:700;margin:4px 0 10px}
.div{height:3px;background:linear-gradient(90deg,#ffd700,transparent);border-radius:2px;margin-bottom:16px;width:60%}
ul{list-style:none;padding:0;margin:0;display:flex;flex-direction:column;gap:8px}
li{color:#ccdeff;font-size:15px;padding-left:20px;position:relative;line-height:1.5}
li::before{content:'▸';color:#ffd700;position:absolute;left:0}
</style></head><body>\n""" + "\n".join(slides_html) + "\n</body></html>"


async def _pptx_upload(owui_url: str, token: str, file_path: str, filename: str) -> Optional[str]:
    """Try to upload via OpenWebUI API. Returns download URL or None."""
    # Try with user token first, then with pipelines key
    tokens_to_try = [t for t in [token, "sk-pipelines-local", "sk-ewgiaPC3A6pPDYHwR8siVA"] if t]
    for tok in tokens_to_try:
        try:
            with open(file_path, "rb") as f:
                data = f.read()
            async with httpx.AsyncClient(timeout=30) as client:
                r = await client.post(
                    f"{owui_url}/api/v1/files/",
                    headers={"Authorization": f"Bearer {tok}"},
                    files={"file": (filename, data,
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
                )
            if r.status_code == 200:
                fid = r.json().get("id") or r.json().get("file_id")
                if fid:
                    return f"{owui_url}/api/v1/files/{fid}/content"
            log.warning("PPTX upload attempt failed: %s", r.status_code)
        except Exception as e:
            log.warning("PPTX upload error: %s", e)
    return None


async def handle_pptx(
    body: dict, base_url: str, api_key: str, model: str,
    owui_url: str, user_info: dict,
    emitter: Optional[Callable] = None,
) -> dict:
    """Full PPTX pipeline: generate → build → upload → reply with HTML preview + link."""
    messages = body.get("messages", [])
    last_text = ""
    for msg in reversed(messages):
        if msg.get("role") == "user":
            c = msg.get("content", "")
            last_text = c if isinstance(c, str) else " ".join(
                p.get("text", "") for p in c if isinstance(p, dict) and p.get("type") == "text"
            )
            break

    topic = _pptx_topic(last_text) or "Презентация"
    num = max(3, min(15, _pptx_slide_count(last_text)))

    async def emit(text, done=False):
        if emitter and callable(emitter):
            await emitter({"type": "status", "data": {"description": text, "done": done}})

    await emit(f"📊 Генерирую презентацию: «{topic[:50]}»…")

    # Web search if topic needs fresh data
    web_context = ""
    needs_search = bool(re.search(
        r"\b(2024|2025|2026|2027|новинк|выход|релиз|анонс|последн|актуальн|свеж"
        r"|new|upcoming|release|latest|recent|announce)\b",
        last_text, re.IGNORECASE,
    ))
    if needs_search:
        await emit(f"🔍 Ищу актуальные данные по теме…")
        try:
            from urllib.parse import quote_plus
            ua = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 Chrome/124.0.0.0 Safari/537.36"
            encoded = quote_plus(topic)
            async with httpx.AsyncClient(timeout=15, follow_redirects=True,
                headers={"User-Agent": ua, "Accept-Language": "ru-RU,ru;q=0.9,en;q=0.8"}) as client:
                r = await client.get(f"https://html.duckduckgo.com/html/?q={encoded}&kl=ru-ru")
            import html as _html
            raw = r.text
            snippet_re = re.compile(r'<a[^>]+class="result__snippet"[^>]*>(.*?)</a>', re.DOTALL)
            link_re = re.compile(r'class="result__a"[^>]+href="([^"]+)"', re.DOTALL)
            snippets = [re.sub(r"<[^>]+>", " ", _html.unescape(s))[:300] for s in snippet_re.findall(raw)[:5]]
            web_context = "\n".join(snippets)
        except Exception as e:
            log.warning("PPTX web search error: %s", e)

    slides_data = await _pptx_generate_slides(base_url, api_key, model, topic, num, web_context)

    await emit("🔨 Собираю .pptx…")
    out_dir = Path("/app/data/pptx")
    out_dir.mkdir(parents=True, exist_ok=True)
    import uuid as _uuid
    filename = f"presentation_{_uuid.uuid4().hex[:8]}.pptx"
    file_path = str(out_dir / filename)
    loop = asyncio.get_event_loop()
    await loop.run_in_executor(None, _pptx_build, topic, slides_data, file_path)

    await emit("📤 Загружаю файл…")
    token = user_info.get("token") or user_info.get("api_key") or api_key
    log.info("PPTX upload token keys available: %s", list(user_info.keys()))
    download_url = await _pptx_upload(owui_url, token, file_path, filename)

    html = _pptx_html_preview(topic, slides_data)
    download_line = (
        f"[⬇️ Скачать {filename}]({download_url})"
        if download_url else
        f"[⬇️ Скачать {filename}](http://localhost/mts/generated_pptx/{filename})"
    )
    reply = f"📊 **Презентация готова** — {len(slides_data) + 1} слайдов\n\n{download_line}\n\n```html\n{html}\n```"

    await emit("✅ Готово", done=True)

    body["messages"] = [
        {"role": "system", "content": "Выведи следующее сообщение пользователю дословно без изменений."},
        {"role": "user", "content": reply},
    ]
    body["model"] = model
    body["stream"] = False
    body["__pptx_handled__"] = True
    return body


# ---------------------------------------------------------------------------
# Filter
# ---------------------------------------------------------------------------

class Filter:
    type = "filter"
    class Valves(BaseModel):
        model_config = {"protected_namespaces": ()}
        pipelines: list[str] = Field(default=["*"], description="Модели к которым применяется фильтр")
        priority: int = Field(default=10, description="Приоритет фильтра (меньше = раньше)")

        base_url: str = "https://api.gpt.mws.ru/v1"
        api_key: str = "sk-ewgiaPC3A6pPDYHwR8siVA"
        openwebui_base_url: str = "http://open-webui:8080"

        model_vision: str = "qwen2.5-vl-72b"
        model_vision_alt: str = "cotype-pro-vl-32b"
        model_audio: str = "whisper-turbo-local"
        model_code: str = "qwen3-coder-480b-a35b"
        model_default: str = "llama-3.3-70b-instruct"
        model_memory: str = "mts-anya"
        model_embed: str = "bge-m3"

        unavailable_models: list[str] = []
        memory_top_k: int = 5
        file_rag_top_k: int = 5
        force_search_on_sources: bool = True

    class UserValves(BaseModel):
        enable_memory: bool = Field(
            default=True,
            description="🧠 Запоминать информацию обо мне между чатами",
        )
        clear_memory: bool = Field(
            default=False,
            description="🗑 Очистить мою память (сбросится после следующего сообщения)",
        )
        show_model_status: bool = Field(
            default=True,
            description="👁 Показывать какая модель выбрана для ответа",
        )

    def __init__(self):
        self.valves = self.Valves()
        # Per-user last-response cache for format conversion: {user_id: text}
        self._last_response: dict = {}

    # ── Helpers ────────────────────────────────────────────────────────────

    def _has_image(self, messages: list) -> bool:
        for msg in messages:
            if isinstance(msg.get("content"), list):
                for part in msg["content"]:
                    if isinstance(part, dict) and part.get("type") == "image_url":
                        return True
        return False

    def _has_audio(self, messages: list) -> bool:
        for msg in messages:
            if isinstance(msg.get("content"), list):
                for part in msg["content"]:
                    if isinstance(part, dict) and part.get("type") in ("audio", "audio_url"):
                        return True
        return False

    def _extract_text(self, messages: list) -> str:
        parts = []
        for msg in messages:
            c = msg.get("content", "")
            if isinstance(c, str):
                parts.append(c)
            elif isinstance(c, list):
                for p in c:
                    if isinstance(p, dict) and p.get("type") == "text":
                        parts.append(p.get("text", ""))
        return " ".join(parts)

    def _is_code_request(self, text: str) -> bool:
        pat = (
            r"\b(code|coding|python|javascript|typescript|java|c\+\+|golang|rust|sql|bash"
            r"|script|function|class|algorithm|debug|refactor|snippet"
            # русские
            r"|напиши\s+код|напишем\s+код|написать\s+код"
            r"|напиши\s+(?:лендинг|сайт|страниц|компонент|функци|скрипт|запрос|апи|бэкенд|фронтенд)"
            r"|сделай\s+(?:лендинг|сайт|страниц|компонент|функци|скрипт)"
            r"|код\s+для|пример\s+кода|html|css|верстк|разработ)\b"
        )
        return bool(re.search(pat, text, re.IGNORECASE)) or "```" in text

    # ── Audio handler ──────────────────────────────────────────────────────

    async def _transcribe_audio(self, messages: list) -> Optional[str]:
        """
        Extract audio data from message parts and transcribe via whisper-turbo-local.
        Supports both audio_url (remote) and inline base64 audio blobs.
        Returns the transcript string, or None on failure.
        """
        for msg in messages:
            if not isinstance(msg.get("content"), list):
                continue
            for part in msg["content"]:
                if not isinstance(part, dict):
                    continue
                ptype = part.get("type")

                if ptype == "audio_url":
                    audio_url = (part.get("audio_url") or {}).get("url", "")
                elif ptype == "audio":
                    audio_url = part.get("url", "") or part.get("data", "")
                else:
                    continue

                if not audio_url:
                    continue

                try:
                    # Decode base64 blob if needed
                    if audio_url.startswith("data:"):
                        header, b64data = audio_url.split(",", 1)
                        mime = header.split(";")[0].replace("data:", "")
                        ext = mime.split("/")[-1] or "wav"
                        audio_bytes = base64.b64decode(b64data)
                        files = {"file": (f"audio.{ext}", audio_bytes, mime)}
                    else:
                        # Remote URL — download first
                        async with httpx.AsyncClient(timeout=30) as client:
                            dl = await client.get(audio_url)
                        dl.raise_for_status()
                        ext = audio_url.rsplit(".", 1)[-1].split("?")[0] or "wav"
                        files = {"file": (f"audio.{ext}", dl.content, f"audio/{ext}")}

                    async with httpx.AsyncClient(timeout=60) as client:
                        resp = await client.post(
                            f"{self.valves.base_url}/audio/transcriptions",
                            headers={"Authorization": f"Bearer {self.valves.api_key}"},
                            data={"model": self.valves.model_audio},
                            files=files,
                        )
                    resp.raise_for_status()
                    return resp.json().get("text", "").strip()
                except Exception as e:
                    log.error("Audio transcription error: %s", e)
                    return None
        return None

    def _normalize_image_parts(self, messages: list) -> list:
        """
        Ensure all image_url parts carry properly formatted base64 data URIs.
        Handles: raw base64 strings, data URIs, and remote URLs (left as-is).
        """
        for msg in messages:
            if not isinstance(msg.get("content"), list):
                continue
            for part in msg["content"]:
                if not isinstance(part, dict) or part.get("type") != "image_url":
                    continue
                img = part.get("image_url") or {}
                url = img.get("url", "")
                # Already a proper data URI or remote URL — nothing to do
                if url.startswith("data:image/") or url.startswith("http"):
                    continue
                # Raw base64 without header — wrap it
                if url and not url.startswith("data:"):
                    # Try to detect format from base64 magic bytes
                    try:
                        raw = base64.b64decode(url[:16])
                        if raw[:4] == b"\x89PNG":
                            mime = "image/png"
                        elif raw[:2] in (b"\xff\xd8",):
                            mime = "image/jpeg"
                        elif raw[:4] == b"RIFF":
                            mime = "image/webp"
                        elif raw[:3] == b"GIF":
                            mime = "image/gif"
                        else:
                            mime = "image/jpeg"
                    except Exception:
                        mime = "image/jpeg"
                    part["image_url"]["url"] = f"data:{mime};base64,{url}"
        return messages

    # ── File RAG handler ───────────────────────────────────────────────────

    def _extract_file_ids(self, messages: list) -> list[str]:
        """
        Collect file IDs from message metadata.
        OpenWebUI attaches files as parts with type='file' or in message metadata.
        """
        file_ids = []
        for msg in messages:
            # Check content parts
            if isinstance(msg.get("content"), list):
                for part in msg["content"]:
                    if isinstance(part, dict) and part.get("type") == "file":
                        fid = part.get("file_id") or part.get("id")
                        if fid:
                            file_ids.append(fid)
            # Check metadata field (OpenWebUI stores file refs here)
            meta = msg.get("metadata") or {}
            for f in meta.get("files", []):
                fid = f.get("id") or f.get("file_id")
                if fid and fid not in file_ids:
                    file_ids.append(fid)
        return file_ids

    async def _fetch_collection_content(
        self, api_token: str, file_ids: list[str], query: str
    ) -> Optional[str]:
        """
        Query OpenWebUI's RAG endpoint for relevant chunks from attached files.
        Uses bge-m3 embeddings for retrieval.
        """
        if not file_ids:
            return None
        try:
            payload = {
                "query": query,
                "collection_names": file_ids,
                "k": self.valves.file_rag_top_k,
                "r": 0.0,
                "hybrid": True,
            }
            headers = {
                "Authorization": f"Bearer {api_token}",
                "Content-Type": "application/json",
            }
            async with httpx.AsyncClient(timeout=20) as client:
                resp = await client.post(
                    f"{self.valves.openwebui_base_url}/api/v1/retrieval/process/query",
                    headers=headers,
                    json=payload,
                )
            if resp.status_code == 404:
                # Fallback: try legacy endpoint
                async with httpx.AsyncClient(timeout=20) as client:
                    resp = await client.post(
                        f"{self.valves.openwebui_base_url}/api/v1/rag/query",
                        headers=headers,
                        json=payload,
                    )
            resp.raise_for_status()
            data = resp.json()
            # Extract text from result documents
            docs = data.get("documents") or data.get("results") or []
            chunks = []
            for doc in docs:
                if isinstance(doc, list):
                    chunks.extend(doc)
                elif isinstance(doc, str):
                    chunks.append(doc)
                elif isinstance(doc, dict):
                    chunks.append(doc.get("content") or doc.get("text") or "")
            return "\n\n".join(c for c in chunks if c).strip() or None
        except Exception as e:
            log.error("File RAG fetch error: %s", e)
            return None

    async def _file_rag_context(
        self, messages: list, query: str, api_token: str
    ) -> Optional[str]:
        """
        Full File RAG pipeline:
        1. Extract file IDs from message metadata
        2. Query OpenWebUI RAG with bge-m3 embeddings
        3. Return formatted [FILE CONTEXT] block or None
        """
        file_ids = self._extract_file_ids(messages)
        if not file_ids:
            return None
        content = await self._fetch_collection_content(api_token, file_ids, query)
        if not content:
            return None
        return f"[FILE CONTEXT]\nИспользуй следующие фрагменты из прикреплённых файлов:\n\n{content}"

    def _inject_file_context(self, messages: list, block: str) -> list:
        """Inject [FILE CONTEXT] block into system prompt."""
        if not block:
            return messages
        for msg in messages:
            if msg.get("role") == "system":
                if "[FILE CONTEXT]" in msg["content"]:
                    msg["content"] = re.sub(
                        r"\[FILE CONTEXT\].*?(?=\n\n|\Z)", block,
                        msg["content"], flags=re.DOTALL,
                    )
                else:
                    msg["content"] = block + "\n\n" + msg["content"]
                return messages
        messages.insert(0, {"role": "system", "content": block})
        return messages

    async def _fetch_file_bytes(self, api_token: str, file_id: str) -> Optional[tuple[str, bytes]]:
        """
        Download raw file bytes from OpenWebUI.
        Returns (filename, bytes) or None on failure.
        """
        headers = {"Authorization": f"Bearer {api_token}"}
        for tok in [api_token, "sk-pipelines-local", self.valves.api_key]:
            try:
                async with httpx.AsyncClient(timeout=30) as client:
                    # Get file metadata first
                    meta_resp = await client.get(
                        f"{self.valves.openwebui_base_url}/api/v1/files/{file_id}",
                        headers={"Authorization": f"Bearer {tok}"},
                    )
                    if meta_resp.status_code != 200:
                        continue
                    meta = meta_resp.json()
                    filename = meta.get("filename") or meta.get("name") or f"{file_id}.bin"
                    # Download content
                    content_resp = await client.get(
                        f"{self.valves.openwebui_base_url}/api/v1/files/{file_id}/content",
                        headers={"Authorization": f"Bearer {tok}"},
                    )
                    if content_resp.status_code == 200:
                        return filename, content_resp.content
            except Exception as e:
                log.warning("File fetch error (id=%s): %s", file_id, e)
        return None

    async def _direct_parse_files(
        self, messages: list, api_token: str
    ) -> Optional[str]:
        """
        Fetch raw file bytes from OpenWebUI and parse them locally using
        file_processor.extract_text_from_any. Handles .pdf/.docx/.txt/.csv/.xlsx.
        Returns a combined [FILE CONTEXT] block or None.
        """
        if not _FILE_PROCESSOR_AVAILABLE:
            return None
        file_ids = self._extract_file_ids(messages)
        if not file_ids:
            return None

        blocks = []
        for fid in file_ids:
            result = await self._fetch_file_bytes(api_token, fid)
            if result is None:
                continue
            filename, data = result
            text = extract_text_from_any(filename, data)
            if text:
                blocks.append(build_file_context_block(filename, text))
                log.info("Parsed file '%s' (%d chars)", filename, len(text))

        if not blocks:
            return None
        return "\n\n---\n\n".join(blocks)

    def _resolve_model(self, preferred: str) -> str:
        return self.valves.model_default if preferred in self.valves.unavailable_models else preferred

    async def _emit(self, emitter, text: str, done: bool = True):
        if emitter and callable(emitter):
            await emitter({"type": "status", "data": {"description": text, "done": done}})

    def _parse_user_valves(self, user_info: dict) -> "Filter.UserValves":
        raw = user_info.get("valves") or {}
        if isinstance(raw, str):
            try:
                raw = json.loads(raw)
            except Exception:
                raw = {}
        elif hasattr(raw, "__dict__"):
            raw = vars(raw)
        elif not isinstance(raw, dict):
            raw = {}
        try:
            return self.UserValves(**raw)
        except Exception:
            return self.UserValves()

    # ── inlet ──────────────────────────────────────────────────────────────

    async def inlet(
        self,
        body: dict,
        __event_emitter__: Optional[Callable[[dict], Awaitable[None]]] = None,
        user: Optional[dict] = None,
        __user__: Optional[dict] = None,
    ) -> dict:
        try:
            messages: list = body.get("messages", [])
            text = self._extract_text(messages)
            user_info = __user__ or user or {}
            user_id = user_info.get("id", "anonymous")
            # Isolate memory per chat session to prevent cross-chat contamination
            chat_id = body.get("chat_id") or body.get("session_id") or ""
            if chat_id:
                user_id = f"{user_id}:{chat_id}"
            # OpenWebUI passes the user's API token for RAG calls
            api_token = user_info.get("token") or self.valves.api_key
            uv = self._parse_user_valves(user_info)

            # ── Если image_gen_pipeline уже обработал запрос — не трогаем ─
            if body.get("__image_handled__"):
                return body

            # ── Если pptx_pipeline уже обработал запрос — не трогаем ──────
            if body.get("__pptx_handled__"):
                return body

            # ── Clear memory ───────────────────────────────────────────────
            if uv.clear_memory:
                _clear_memory(user_id)
                await self._emit(__event_emitter__, "🗑 Память очищена")

            # ── Memory pipeline ────────────────────────────────────────────
            if uv.enable_memory:
                user_msgs = [m for m in messages if m.get("role") == "user"]
                if user_msgs:
                    last_content = user_msgs[-1].get("content", "")
                    latest = last_content if isinstance(last_content, str) else self._extract_text([user_msgs[-1]])
                    new_facts = _extract_facts(latest)
                    if new_facts:
                        for key, value in new_facts.items():
                            _save_fact(user_id, key, value)
                        log.info("Fast-extracted facts for '%s': %s", user_id, new_facts)

                relevant = await _semantic_search(
                    self.valves.base_url, self.valves.api_key,
                    user_id, text, top_k=self.valves.memory_top_k,
                )
                if relevant:
                    block = _build_context_block(relevant)
                    messages = _inject_context(messages, block)
                    body["messages"] = messages
                    names = ", ".join(f"{f['key']}={f['value']}" for f in relevant)
                    await self._emit(__event_emitter__, f"🧠 Память: {names}")

            # ── Audio: transcribe first, then route to text model ──────────
            if self._has_audio(messages):
                await self._emit(__event_emitter__, "🎙 Транскрибирую аудио...", done=False)
                transcript = await self._transcribe_audio(messages)
                if transcript:
                    # Replace audio parts with the transcript text
                    for msg in messages:
                        if not isinstance(msg.get("content"), list):
                            continue
                        text_parts = [
                            p for p in msg["content"]
                            if isinstance(p, dict) and p.get("type") == "text"
                        ]
                        audio_parts = [
                            p for p in msg["content"]
                            if isinstance(p, dict) and p.get("type") in ("audio", "audio_url")
                        ]
                        if audio_parts:
                            text_parts.append({"type": "text", "text": f"[Транскрипция]: {transcript}"})
                            msg["content"] = text_parts
                    body["messages"] = messages
                    text = self._extract_text(messages)  # refresh for downstream detection
                    await self._emit(__event_emitter__, f"🎙 Транскрипция: {transcript[:80]}{'…' if len(transcript) > 80 else ''}")

                    # После транскрипции — проверяем, не запрос ли это на генерацию картинки.
                    # Если да — не трогаем модель, image_gen_pipeline сам перехватит.
                    if _is_image_gen_request(text):
                        chosen, label = self._resolve_model(self.valves.model_default), "🎙🎨 Аудио→Картинка"
                    elif self._is_code_request(text):
                        chosen, label = self._resolve_model(self.valves.model_code), "💻 Аудио→Код"
                    else:
                        chosen, label = self._resolve_model(self.valves.model_default), "🎙 Аудио→Текст"
                else:
                    chosen, label = self._resolve_model(self.valves.model_audio), "🎙 Аудио"

            # ── Vision: normalize base64 images, pick VLM ─────────────────
            elif self._has_image(messages):
                messages = self._normalize_image_parts(messages)
                body["messages"] = messages
                # Prefer qwen2.5-vl; fall back to cotype-pro-vl-32b if unavailable
                preferred_vl = self.valves.model_vision
                if preferred_vl in self.valves.unavailable_models:
                    preferred_vl = self.valves.model_vision_alt
                chosen = self._resolve_model(preferred_vl)
                label = "👁 Изображение"
                vision_sys = (
                    "Ты — ассистент для анализа изображений. "
                    "Отвечай ТОЛЬКО на русском языке. "
                    "Описывай изображение чётко и конкретно. "
                    "Никаких иероглифов, случайных символов или кода."
                )
                has_sys = any(m.get("role") == "system" for m in messages)
                if has_sys:
                    for m in messages:
                        if m.get("role") == "system":
                            m["content"] = vision_sys + "\n\n" + m["content"]
                            break
                else:
                    messages.insert(0, {"role": "system", "content": vision_sys})
                body["messages"] = messages

            # ── File RAG: retrieve context, inject, route to default/code ──
            elif self._extract_file_ids(messages):
                await self._emit(__event_emitter__, "📄 Обрабатываю файлы...", done=False)
                user_msgs = [m for m in messages if m.get("role") == "user"]
                query = self._extract_text(user_msgs[-1:]) if user_msgs else text

                # 1. Try direct local parsing (.pdf/.docx/.txt/.csv/.xlsx)
                file_block = await self._direct_parse_files(messages, api_token)
                if file_block:
                    messages = self._inject_file_context(messages, file_block)
                    body["messages"] = messages
                    await self._emit(__event_emitter__, "📄 Файл разобран и добавлен в контекст")
                else:
                    # 2. Fallback: OpenWebUI RAG collection query
                    await self._emit(__event_emitter__, "📄 Поиск по файлам (bge-m3)...", done=False)
                    file_block = await self._file_rag_context(messages, query, api_token)
                    if file_block:
                        messages = self._inject_file_context(messages, file_block)
                        body["messages"] = messages
                        await self._emit(__event_emitter__, "📄 Контекст из файлов добавлен")
                    else:
                        await self._emit(__event_emitter__, "📄 Файлы прикреплены (контекст не найден)")

                if self._is_code_request(text) and not _is_image_gen_request(text):
                    chosen, label = self._resolve_model(self.valves.model_code), "💻 Файл→Код"
                else:
                    chosen, label = self._resolve_model(self.valves.model_default), "📄 Файл→Текст"

            # ── Code-edit: "добавь логотип в код выше" → code model, не image ─
            elif _is_code_edit_request(text):
                chosen, label = self._resolve_model(self.valves.model_code), "✏️ Правка кода"

            elif self._is_code_request(text) and not _is_image_gen_request(text):
                chosen, label = self._resolve_model(self.valves.model_code), "💻 Код"
            else:
                chosen, label = self._resolve_model(self.valves.model_default), "💬 Текст"

            # ── Презентация — не переключаем модель, pptx_pipeline сам обработает ─
            if _is_pptx_request(text):
                return body

            # ── File generation / conversion — inject hidden tag instruction ──
            # Covers both: "сохрани как PDF", "конвертируй в Word", "сделай файл"
            _is_file_cmd = (
                (_FILE_GEN_RE.search(text) or _CONVERT_RE.search(text))
                and not _is_pptx_request(text)
            )
            if _is_file_cmd and _FILE_GENERATOR_AVAILABLE:
                fmt = _detect_file_format(text)
                if fmt:
                    cached = self._last_response.get(user_id, "")
                    if cached:
                        # We already have content — generate immediately, no LLM needed
                        await self._emit(__event_emitter__, f"📄 Создаю {fmt.upper()}…", done=False)
                        result = generate_file(cached, fmt)
                        if result:
                            _, filename = result
                            link = make_download_link(filename, self.valves.openwebui_base_url)
                            body["messages"] = [
                                {"role": "system", "content": "Выведи следующее сообщение пользователю дословно без изменений."},
                                {"role": "user", "content": f"✅ Готово!\n\n{link}"},
                            ]
                            body["model"] = self.valves.model_default
                            body["stream"] = False
                            body["__file_handled__"] = True
                            await self._emit(__event_emitter__, f"✅ {fmt.upper()} готов", done=True)
                            return body
                    # No cached content — instruct LLM to generate content AND append the tag
                    tag_instruction = (
                        f"[SYSTEM: FILE GENERATION REQUIRED]\n"
                        f"The user wants a {fmt.upper()} file. Follow these steps EXACTLY:\n"
                        f"1. Write the complete document content (everything that should be in the file).\n"
                        f"2. On the very last line of your response, append EXACTLY this tag and nothing after it:\n"
                        f"[ACTION:GENERATE_FILE(type={fmt.upper()})]\n"
                        f"CRITICAL: Do NOT say 'I will create', do NOT explain. "
                        f"Just write the content, then the tag on its own line. "
                        f"The system intercepts this tag and builds the real file automatically."
                    )
                    for msg in messages:
                        if msg.get("role") == "system":
                            msg["content"] = tag_instruction + "\n\n" + msg["content"]
                            break
                    else:
                        messages.insert(0, {"role": "system", "content": tag_instruction})
                    body["messages"] = messages
                    body["__pending_file_fmt__"] = fmt
                    log.info("File tag instruction injected for format: %s", fmt)

            # ── Force web search for factual / sourced queries ────────────
            if self.valves.force_search_on_sources and _needs_web_search(text):
                citation_block = (
                    "[CITATION REQUIRED]\n"
                    "Пользователь запросил источники или актуальные данные. "
                    "Используй блок [WEB SEARCH RESULTS] для ответа. "
                    "Для КАЖДОГО факта или утверждения включи реальную кликабельную ссылку "
                    "в формате [название источника](https://полный-url). "
                    "НЕ придумывай URL. Если реальных ссылок нет — честно скажи об этом."
                )
                messages = _inject_context(messages, citation_block)
                body["messages"] = messages
                log.info("Citation instruction injected for user '%s'", user_id)

            body["model"] = chosen
            if uv.show_model_status:
                await self._emit(__event_emitter__, f"{label} → {chosen}")

        except Exception as exc:
            log.error("Filter inlet error: %s", exc)
            await self._emit(__event_emitter__, f"⚠️ Ошибка фильтра: {exc}")
            body["model"] = self.valves.model_default

        return body

    # ── outlet ─────────────────────────────────────────────────────────────

    async def outlet(
        self,
        body: dict,
        __event_emitter__: Optional[Callable[[dict], Awaitable[None]]] = None,
        user: Optional[dict] = None,
        __user__: Optional[dict] = None,
    ) -> dict:
        """
        Post-processing:
        1. Intercept [ACTION:GENERATE_FILE(type=EXT)] tag → build file → replace with download link
        2. Cache last assistant response for future format conversion
        3. Run memory summarizer in background
        """
        try:
            user_info = __user__ or user or {}
            user_id = user_info.get("id", "anonymous")
            # Isolate memory per chat session
            chat_id = body.get("chat_id") or body.get("session_id") or ""
            if chat_id:
                user_id = f"{user_id}:{chat_id}"
            uv = self._parse_user_valves(user_info)
            messages = body.get("messages", [])

            # ── Find last assistant message ───────────────────────────────
            last_assistant_idx = None
            for i in reversed(range(len(messages))):
                if messages[i].get("role") == "assistant":
                    last_assistant_idx = i
                    break

            if last_assistant_idx is not None:
                content = messages[last_assistant_idx].get("content", "")
                if isinstance(content, str) and content.strip():

                    # ── Detect action tag (new format) ────────────────────
                    tag_match = _ACTION_TAG_RE.search(content)
                    detected_fmt = tag_match.group(1).lower() if tag_match else None

                    # Strip the tag from content
                    clean_content = _ACTION_TAG_RE.sub("", content).rstrip()

                    # ── Cache clean response for future conversions ───────
                    self._last_response[user_id] = clean_content.strip()

                    # ── Generate file if tag found or pending ─────────────
                    fmt = detected_fmt or body.get("__pending_file_fmt__")
                    if fmt and _FILE_GENERATOR_AVAILABLE and clean_content.strip():
                        fmt_normalized = {"word": "docx", "text": "txt",
                                          "excel": "xlsx"}.get(fmt, fmt)
                        result = generate_file(clean_content.strip(), fmt_normalized)
                        if result:
                            _, filename = result
                            link = make_download_link(filename, self.valves.openwebui_base_url)
                            messages[last_assistant_idx]["content"] = (
                                clean_content.rstrip() + f"\n\n{link}"
                            )
                            body["messages"] = messages
                            log.info("File generated in outlet: %s", filename)
                        else:
                            log.error("File generation failed for format: %s", fmt)
                    elif tag_match:
                        # Tag found but generator unavailable — at least strip the tag
                        messages[last_assistant_idx]["content"] = clean_content
                        body["messages"] = messages

            # ── Memory summarizer ─────────────────────────────────────────
            if uv.enable_memory:
                asyncio.create_task(
                    _summarize_and_save(
                        self.valves.base_url,
                        self.valves.api_key,
                        self.valves.model_memory,
                        user_id,
                        messages,
                    )
                )
        except Exception as e:
            log.error("Outlet error: %s", e)
        return body


# OpenWebUI Pipelines compatibility alias
Pipeline = Filter
