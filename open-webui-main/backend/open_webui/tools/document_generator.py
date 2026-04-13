"""
Document Generator Tool
=======================
Generates real .docx, .xlsx, .pptx files from model-provided content
and returns a download link directly in the chat.

Requirements (auto-installed via frontmatter):
  python-docx, openpyxl, python-pptx
"""

# requirements: python-docx, openpyxl, python-pptx

import io
import json
import logging
import os
import uuid
import time
from typing import Optional, Callable, Awaitable

log = logging.getLogger(__name__)


class Tools:
    def __init__(self):
        pass

    async def create_word_document(
        self,
        title: str,
        content: str,
        __user__: dict = {},
        __request__=None,
        __event_emitter__: Optional[Callable[[dict], Awaitable[None]]] = None,
    ) -> str:
        """
        Create a Word (.docx) document and return a download link.
        Use this when the user asks to create a Word document, report, or .docx file.

        :param title: Document title (used as filename and heading)
        :param content: Full document content in plain text or markdown.
                        Use '# Heading', '## Subheading', '- bullet' formatting.
        """
        try:
            from docx import Document
            from docx.shared import Pt, RGBColor
            from docx.enum.text import WD_ALIGN_PARAGRAPH

            doc = Document()

            # Title
            heading = doc.add_heading(title, level=0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

            # Parse content line by line
            for line in content.split('\n'):
                stripped = line.strip()
                if not stripped:
                    doc.add_paragraph('')
                elif stripped.startswith('### '):
                    doc.add_heading(stripped[4:], level=3)
                elif stripped.startswith('## '):
                    doc.add_heading(stripped[3:], level=2)
                elif stripped.startswith('# '):
                    doc.add_heading(stripped[2:], level=1)
                elif stripped.startswith('- ') or stripped.startswith('* '):
                    p = doc.add_paragraph(stripped[2:], style='List Bullet')
                elif stripped.startswith('1. ') or (len(stripped) > 2 and stripped[0].isdigit() and stripped[1] == '.'):
                    text = stripped.split('. ', 1)[1] if '. ' in stripped else stripped
                    doc.add_paragraph(text, style='List Number')
                elif stripped.startswith('**') and stripped.endswith('**'):
                    p = doc.add_paragraph()
                    run = p.add_run(stripped.strip('*'))
                    run.bold = True
                else:
                    doc.add_paragraph(stripped)

            buf = io.BytesIO()
            doc.save(buf)
            buf.seek(0)

            meta = await _upload_file(
                buf, f'{_safe_filename(title)}.docx',
                'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                __user__, __request__, __event_emitter__
            )
            return json.dumps(
                {
                    'status': 'ok',
                    'url': meta['url'],
                    'filename': meta['filename'],
                    'file_id': meta['file_id'],
                    'size': meta['size'],
                    'content_type': meta['content_type'],
                }
            )

        except Exception as e:
            log.exception(e)
            return json.dumps({'status': 'error', 'message': str(e)})

    async def create_excel_spreadsheet(
        self,
        title: str,
        data: str,
        __user__: dict = {},
        __request__=None,
        __event_emitter__: Optional[Callable[[dict], Awaitable[None]]] = None,
    ) -> str:
        """
        Create an Excel (.xlsx) spreadsheet and return a download link.
        Use this when the user asks to create a table, spreadsheet, or .xlsx file.

        :param title: Spreadsheet title (used as filename and sheet name)
        :param data: Table data as JSON array of arrays, e.g.:
                     [["Name","Age","City"],["Alice",30,"Moscow"],["Bob",25,"SPb"]]
                     OR as CSV-like text with rows separated by newlines and columns by '|' or ','.
        """
        try:
            import openpyxl
            from openpyxl.styles import Font, PatternFill, Alignment

            wb = openpyxl.Workbook()
            ws = wb.active
            ws.title = title[:31]  # Excel sheet name limit

            # Parse data
            rows = _parse_table_data(data)

            header_fill = PatternFill(start_color='366092', end_color='366092', fill_type='solid')
            header_font = Font(color='FFFFFF', bold=True)

            for r_idx, row in enumerate(rows, start=1):
                for c_idx, cell_val in enumerate(row, start=1):
                    cell = ws.cell(row=r_idx, column=c_idx, value=cell_val)
                    if r_idx == 1:
                        cell.fill = header_fill
                        cell.font = header_font
                        cell.alignment = Alignment(horizontal='center')

            # Auto-width columns
            for col in ws.columns:
                max_len = max((len(str(c.value or '')) for c in col), default=10)
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 4, 50)

            buf = io.BytesIO()
            wb.save(buf)
            buf.seek(0)

            meta = await _upload_file(
                buf, f'{_safe_filename(title)}.xlsx',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                __user__, __request__, __event_emitter__
            )
            return json.dumps(
                {
                    'status': 'ok',
                    'url': meta['url'],
                    'filename': meta['filename'],
                    'file_id': meta['file_id'],
                    'size': meta['size'],
                    'content_type': meta['content_type'],
                }
            )

        except Exception as e:
            log.exception(e)
            return json.dumps({'status': 'error', 'message': str(e)})

    async def create_presentation(
        self,
        title: str,
        slides: str,
        __user__: dict = {},
        __request__=None,
        __event_emitter__: Optional[Callable[[dict], Awaitable[None]]] = None,
    ) -> str:
        """
        Create a PowerPoint (.pptx) presentation and return a download link.
        Use this when the user asks to create a presentation, slides, or .pptx file.

        :param title: Presentation title (used as filename and title slide)
        :param slides: Slides as JSON array, each item: {"title": "...", "bullets": ["...", "..."]}
                       OR as plain text with slides separated by '---', title on first line,
                       bullets as '- item' lines.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            slide_data = _parse_slides_data(slides, title)

            for i, slide_info in enumerate(slide_data):
                if i == 0:
                    layout = prs.slide_layouts[0]  # Title slide
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = slide_info.get('title', title)
                    if len(slide.placeholders) > 1:
                        slide.placeholders[1].text = slide_info.get('subtitle', '')
                else:
                    layout = prs.slide_layouts[1]  # Title + Content
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = slide_info.get('title', f'Slide {i + 1}')
                    tf = slide.placeholders[1].text_frame
                    tf.clear()
                    bullets = slide_info.get('bullets', [])
                    for j, bullet in enumerate(bullets):
                        if j == 0:
                            tf.paragraphs[0].text = bullet
                            tf.paragraphs[0].level = 0
                        else:
                            p = tf.add_paragraph()
                            p.text = bullet
                            p.level = 0

            buf = io.BytesIO()
            prs.save(buf)
            buf.seek(0)

            meta = await _upload_file(
                buf, f'{_safe_filename(title)}.pptx',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                __user__, __request__, __event_emitter__
            )
            return json.dumps(
                {
                    'status': 'ok',
                    'url': meta['url'],
                    'filename': meta['filename'],
                    'file_id': meta['file_id'],
                    'size': meta['size'],
                    'content_type': meta['content_type'],
                }
            )

        except Exception as e:
            log.exception(e)
            return json.dumps({'status': 'error', 'message': str(e)})


# ── Helpers ───────────────────────────────────────────────────────────────────

def _safe_filename(name: str) -> str:
    return ''.join(c if c.isalnum() or c in ' _-' else '_' for c in name).strip()[:60]


def _parse_table_data(data: str) -> list[list]:
    """Parse JSON array-of-arrays or pipe/comma-separated text into rows."""
    data = data.strip()
    try:
        parsed = json.loads(data)
        if isinstance(parsed, list):
            return [[str(cell) for cell in row] for row in parsed]
    except (json.JSONDecodeError, TypeError):
        pass

    rows = []
    for line in data.split('\n'):
        line = line.strip()
        if not line or line.startswith('---'):
            continue
        # Remove markdown table separators like |---|---|
        if all(c in '-|: ' for c in line):
            continue
        sep = '|' if '|' in line else ','
        cells = [c.strip() for c in line.split(sep)]
        cells = [c for c in cells if c]  # remove empty from leading/trailing |
        if cells:
            rows.append(cells)
    return rows


def _parse_slides_data(data: str, default_title: str) -> list[dict]:
    """Parse JSON array or '---' separated text into slide dicts."""
    data = data.strip()
    try:
        parsed = json.loads(data)
        if isinstance(parsed, list):
            return parsed
    except (json.JSONDecodeError, TypeError):
        pass

    slides = []
    current = None
    for line in data.split('\n'):
        stripped = line.strip()
        if stripped == '---':
            if current is not None:
                slides.append(current)
            current = {'title': '', 'bullets': []}
        elif current is None:
            current = {'title': default_title, 'bullets': []}
            if stripped:
                current['title'] = stripped
        elif not current['title'] and stripped and not stripped.startswith('-'):
            current['title'] = stripped
        elif stripped.startswith('- ') or stripped.startswith('* '):
            current['bullets'].append(stripped[2:])
        elif stripped and not stripped.startswith('#'):
            current['bullets'].append(stripped)

    if current is not None:
        slides.append(current)

    if not slides:
        slides = [{'title': default_title, 'bullets': [], 'subtitle': ''}]

    return slides


async def _upload_file(
    buf: io.BytesIO,
    filename: str,
    content_type: str,
    user: dict,
    request,
    event_emitter,
) -> dict:
    """Upload bytes to the file store and return URL and metadata for downloads."""
    from open_webui.storage.provider import Storage
    from open_webui.models.files import Files, FileForm
    from fastapi import UploadFile
    import uuid as _uuid

    file_id = str(_uuid.uuid4())
    storage_filename = f'{file_id}_{filename}'

    contents, file_path = Storage.upload_file(buf, storage_filename, {
        'OpenWebUI-File-Id': file_id,
        'OpenWebUI-User-Id': user.get('id', 'system'),
    })

    file_item = Files.insert_new_file(
        user.get('id', 'system'),
        FileForm(**{
            'id': file_id,
            'filename': filename,
            'path': file_path,
            'data': {},
            'meta': {
                'name': filename,
                'content_type': content_type,
                'size': len(contents),
            },
        }),
    )

    if event_emitter:
        await event_emitter({
            'type': 'status',
            'data': {
                'description': f'📎 {filename} готов к скачиванию',
                'done': True,
            },
        })

    # Return URL that the frontend can render as a download link
    base_url = ''
    if request:
        try:
            base_url = str(request.base_url).rstrip('/')
        except Exception:
            pass

    download_url = f'{base_url}/api/v1/files/{file_id}/content?attachment=true'
    return {
        'url': download_url,
        'file_id': file_id,
        'filename': filename,
        'size': len(contents),
        'content_type': content_type,
    }
