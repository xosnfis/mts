"""
title: PPTX Generator
id: pptx_pipeline
description: Генерация презентаций .pptx. Скажи "сделай презентацию по [теме]" — получишь готовый файл со слайдами.
author: GPTHub
version: 1.0.0
"""

import re
import json
import logging
import uuid
import asyncio
from pathlib import Path
from typing import Optional, Callable, Awaitable

import httpx
from pydantic import BaseModel, Field

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Trigger & parsing
# ---------------------------------------------------------------------------

_PPTX_RE = re.compile(
    r"\bпрезентаци[юяи]\b|\bслайды\b|\bpptx\b|\bpowerpoint\b"
    r"|presentation|make\s+slides|create\s+slides",
    re.IGNORECASE,
)

_SLIDE_COUNT_RE = re.compile(
    r"\b(один|два|три|четыре|пять|шесть|семь|восемь|девять|десять|\d+)\s+слайд\w*\b"
    r"|\bиз\s+(один|два|три|четыре|пять|шесть|семь|восемь|девять|десять|\d+)\s+слайд\w*\b"
    r"|\bтоп[\s\-]?(один|два|три|четыре|пять|шесть|семь|восемь|девять|десять|\d+)\b"
    r"|\b(один|два|три|четыре|пять|шесть|семь|восемь|девять|десять|\d+)\s+(?:полезн|лучш|топ|интересн|основн|главн|важн)\w*\b",
    re.IGNORECASE,
)

_WORD_TO_NUM = {
    "один": 1, "два": 2, "три": 3, "четыре": 4, "пять": 5,
    "шесть": 6, "семь": 7, "восемь": 8, "девять": 9, "десять": 10,
}


def _slide_count(text: str) -> int:
    m = _SLIDE_COUNT_RE.search(text)
    if not m:
        return 7
    val = next((g for g in m.groups() if g is not None), None)
    if not val:
        return 7
    val = val.lower()
    return max(3, min(15, int(val) if val.isdigit() else _WORD_TO_NUM.get(val, 7)))


def _topic(text: str) -> str:
    m = re.search(r"на\s+тему\s+[«\"]?(.+?)[»\"]?(?:\.|$)", text, re.IGNORECASE)
    if m:
        return m.group(1).strip(" .,»\"")
    m = re.search(r"(?:презентацию?\s+(?:про|о|об|по)|слайды\s+(?:про|о|об|по))\s+(.+?)(?:\.|$)", text, re.IGNORECASE)
    if m:
        return m.group(1).strip(" .,")
    cleaned = re.sub(
        r"\b(сделай|создай|напиши|сгенерируй|подготовь|составь|сгенерирую|make|create|generate|build)"
        r"\s+(?:мне\s+)?(?:a\s+)?(?:презентацию?|слайды|presentation|pptx|powerpoint|slides)\b",
        "", text, flags=re.IGNORECASE,
    )
    cleaned = re.sub(r"\b(\d+|один|два|три|четыре|пять|шесть|семь|восемь|девять|десять)\s+слайд\w*\b", "", cleaned, flags=re.IGNORECASE)
    return re.sub(r"\s{2,}", " ", cleaned).strip(" ,.") or text


# ---------------------------------------------------------------------------
# LLM slide content
# ---------------------------------------------------------------------------

async def _gen_slides(base_url: str, api_key: str, model: str, topic: str, num: int) -> list[dict]:
    prompt = (
        f"Создай содержимое для презентации на тему: «{topic}».\n"
        f"Нужно ровно {num} слайдов (не считая титульного).\n"
        f"Верни ТОЛЬКО валидный JSON-массив без пояснений, без markdown-блоков. Формат:\n"
        f'[{{"title":"Заголовок слайда","bullets":["Пункт 1","Пункт 2","Пункт 3"]}}]\n'
        f"Требования: 3-5 пунктов на слайд, каждый до 15 слов. Язык: русский.\n"
        f"ВАЖНО: верни ТОЛЬКО JSON-массив, начиная с [ и заканчивая ]."
    )
    try:
        async with httpx.AsyncClient(timeout=60) as client:
            r = await client.post(
                f"{base_url}/chat/completions",
                headers={"Authorization": f"Bearer {api_key}"},
                json={"model": model, "messages": [{"role": "user", "content": prompt}],
                      "temperature": 0.3, "max_tokens": 2000},
            )
        r.raise_for_status()
        raw = r.json()["choices"][0]["message"]["content"].strip()

        # Try direct JSON parse first
        try:
            data = json.loads(raw)
            if isinstance(data, list):
                return _normalize_slides(data, num)
        except json.JSONDecodeError:
            pass

        # Try extracting JSON array from text
        m = re.search(r"\[.*\]", raw, re.DOTALL)
        if m:
            try:
                data = json.loads(m.group(0))
                if isinstance(data, list):
                    return _normalize_slides(data, num)
            except json.JSONDecodeError:
                pass

        # Fallback: parse the text response directly into slides
        log.warning("JSON parse failed, falling back to text parsing for slides")
        return _parse_text_to_slides(raw, num)

    except Exception as e:
        log.error("Slide gen error: %s", e)
    return [{"title": f"Раздел {i}", "bullets": [f"Содержимое раздела {i}"]} for i in range(1, num + 1)]


def _normalize_slides(data: list, num: int) -> list[dict]:
    """Ensure each slide has title and bullets list."""
    result = []
    for item in data[:num]:
        if not isinstance(item, dict):
            continue
        title = str(item.get("title") or item.get("заголовок") or "Слайд")
        bullets = item.get("bullets") or item.get("пункты") or item.get("content") or []
        if isinstance(bullets, str):
            bullets = [b.strip() for b in bullets.split("\n") if b.strip()]
        result.append({"title": title, "bullets": [str(b) for b in bullets if b]})
    # Pad if needed
    while len(result) < num:
        result.append({"title": f"Раздел {len(result)+1}", "bullets": ["Содержимое"]})
    return result


def _parse_text_to_slides(text: str, num: int) -> list[dict]:
    """
    Parse free-form LLM text into slides.
    Looks for patterns like 'Слайд N:', '## Title', bold headers, etc.
    """
    slides = []
    current_title = ""
    current_bullets = []

    slide_header_re = re.compile(
        r"^(?:слайд\s*\d+\s*[:\-–]?\s*|##?\s*|заголовок\s*[:\-]?\s*)(.+)$",
        re.IGNORECASE,
    )
    bullet_re = re.compile(r"^[\-\*•▸]\s+(.+)$")

    for line in text.split("\n"):
        line = line.strip()
        if not line:
            continue

        header_m = slide_header_re.match(line)
        bullet_m = bullet_re.match(line)

        if header_m and len(line) < 80:
            # Save previous slide
            if current_title and current_bullets:
                slides.append({"title": current_title, "bullets": current_bullets})
            current_title = header_m.group(1).strip(" *:#")
            current_bullets = []
        elif bullet_m:
            current_bullets.append(bullet_m.group(1).strip())
        elif current_title and line and not line.startswith("["):
            # Treat as bullet if we have a title
            if len(line) < 120:
                current_bullets.append(line.strip(" *"))

    # Save last slide
    if current_title and current_bullets:
        slides.append({"title": current_title, "bullets": current_bullets})

    if not slides:
        # Last resort: split into chunks
        lines = [l.strip() for l in text.split("\n") if l.strip() and len(l.strip()) > 5]
        chunk = max(1, len(lines) // num)
        for i in range(num):
            chunk_lines = lines[i*chunk:(i+1)*chunk]
            slides.append({
                "title": chunk_lines[0][:60] if chunk_lines else f"Раздел {i+1}",
                "bullets": chunk_lines[1:6] if len(chunk_lines) > 1 else [f"Содержимое {i+1}"],
            })

    return slides[:num] or [{"title": f"Раздел {i}", "bullets": ["Содержимое"]} for i in range(1, num+1)]


# ---------------------------------------------------------------------------
# PPTX builder
# ---------------------------------------------------------------------------

def _build_pptx(topic: str, slides: list[dict], path: str) -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    BG = RGBColor(0x1A, 0x1A, 0x2E)
    HL = RGBColor(0x0F, 0x3C, 0x78)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xCC, 0xDD, 0xFF)
    GOLD = RGBColor(0xFF, 0xD7, 0x00)

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    def bg(slide, c):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = c

    def tb(slide, text, l, t, w, h, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color

    def rect(slide, l, t, w, h, c):
        s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        s.fill.solid()
        s.fill.fore_color.rgb = c
        s.line.fill.background()

    # Title slide
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg(sl, BG)
    rect(sl, 0, 2.8, 13.33, 0.06, GOLD)
    tb(sl, topic, 1.0, 1.2, 11.33, 1.8, size=40, bold=True, align=PP_ALIGN.CENTER)
    tb(sl, "Выполнил/а: _______________", 1.0, 3.2, 11.33, 0.5, size=22, color=LIGHT, align=PP_ALIGN.CENTER)
    tb(sl, "Проверил/а: _______________", 1.0, 3.75, 11.33, 0.5, size=22, color=LIGHT, align=PP_ALIGN.CENTER)

    for sd in slides:
        sl = prs.slides.add_slide(prs.slide_layouts[6])
        bg(sl, BG)
        rect(sl, 0, 0, 13.33, 1.3, HL)
        tb(sl, sd.get("title", ""), 0.4, 0.15, 12.5, 1.0, size=28, bold=True)
        rect(sl, 0.4, 1.45, 12.5, 0.04, GOLD)
        y = 1.65
        for b in sd.get("bullets", []):
            tb(sl, f"▸  {b}", 0.6, y, 12.0, 0.55, size=18, color=LIGHT)
            y += 0.65

    prs.save(path)


# ---------------------------------------------------------------------------
# HTML preview
# ---------------------------------------------------------------------------

def _html_preview(topic: str, slides: list[dict]) -> str:
    parts = [f"""  <div class="slide ts">
    <div class="sn">1</div><h1>{topic}</h1>
    <div class="meta"><span>Выполнил/а: _______________</span><span>Проверил/а: _______________</span></div>
  </div>"""]
    for i, s in enumerate(slides, 2):
        bl = "\n".join(f"    <li>{b}</li>" for b in s.get("bullets", []))
        parts.append(f"""  <div class="slide cs">
    <div class="sn">{i}</div><h2>{s.get('title','')}</h2>
    <div class="dv"></div><ul>{bl}</ul>
  </div>""")
    return (
        '<!DOCTYPE html><html lang="ru"><head><meta charset="UTF-8"><style>'
        'body{font-family:Segoe UI,Arial,sans-serif;background:#0d0d1a;margin:0;padding:16px;display:flex;flex-direction:column;gap:16px}'
        '.slide{max-width:860px;border-radius:10px;position:relative;overflow:hidden;box-shadow:0 4px 20px rgba(0,0,0,.5)}'
        '.ts{background:linear-gradient(135deg,#1a1a2e,#16213e,#0f3c78);min-height:240px;padding:32px 40px}'
        '.cs{background:linear-gradient(160deg,#1a1a2e,#16213e);padding:32px 40px 28px}'
        '.sn{position:absolute;top:12px;right:16px;color:rgba(255,215,0,.5);font-size:12px;font-weight:600}'
        'h1{color:#fff;font-size:28px;font-weight:700;margin:20px 0 24px}'
        '.meta{display:flex;gap:32px;color:#aabbdd;font-size:14px;border-top:2px solid #ffd700;padding-top:14px}'
        'h2{color:#fff;font-size:20px;font-weight:700;margin:4px 0 10px}'
        '.dv{height:3px;background:linear-gradient(90deg,#ffd700,transparent);border-radius:2px;margin-bottom:16px;width:60%}'
        'ul{list-style:none;padding:0;margin:0;display:flex;flex-direction:column;gap:8px}'
        'li{color:#ccdeff;font-size:15px;padding-left:20px;position:relative;line-height:1.5}'
        "li::before{content:'▸';color:#ffd700;position:absolute;left:0}"
        '</style></head><body>\n'
        + "\n".join(parts)
        + '\n</body></html>'
    )


# ---------------------------------------------------------------------------
# Pipeline
# ---------------------------------------------------------------------------

class Pipeline:
    type = "filter"

    class Valves(BaseModel):
        model_config = {"protected_namespaces": ()}
        pipelines: list[str] = Field(default=["*"])
        priority: int = Field(default=0)
        base_url: str = Field(default="https://api.gpt.mws.ru/v1")
        api_key: str = Field(default="sk-ewgiaPC3A6pPDYHwR8siVA")
        model: str = Field(default="llama-3.3-70b-instruct")
        output_dir: str = Field(default="/app/data/pptx")
        download_base_url: str = Field(default="http://localhost/mts/generated_pptx")

    class UserValves(BaseModel):
        enable_pptx: bool = Field(default=True, description="📊 Генерация презентаций (.pptx)")

    def __init__(self):
        self.valves = self.Valves()

    def _parse_uv(self, user_info: dict) -> "Pipeline.UserValves":
        raw = user_info.get("valves") or {}
        if isinstance(raw, str):
            try:
                raw = json.loads(raw)
            except Exception:
                raw = {}
        elif hasattr(raw, "__dict__"):
            raw = vars(raw)
        elif not isinstance(raw, dict):
            raw = {}
        try:
            return self.UserValves(**raw)
        except Exception:
            return self.UserValves()

    def _last_text(self, messages: list) -> str:
        for msg in reversed(messages):
            if msg.get("role") == "user":
                c = msg.get("content", "")
                if isinstance(c, str):
                    return c
                if isinstance(c, list):
                    return " ".join(p.get("text", "") for p in c
                                    if isinstance(p, dict) and p.get("type") == "text")
        return ""

    async def _emit(self, emitter, text: str, done: bool = False):
        if emitter and callable(emitter):
            await emitter({"type": "status", "data": {"description": text, "done": done}})

    async def inlet(
        self,
        body: dict,
        __event_emitter__: Optional[Callable[[dict], Awaitable[None]]] = None,
        user: Optional[dict] = None,
        __user__: Optional[dict] = None,
    ) -> dict:
        try:
            user_info = __user__ or user or {}
            uv = self._parse_uv(user_info)
            if not uv.enable_pptx:
                return body

            messages = body.get("messages", [])
            last_text = self._last_text(messages)

            if not last_text and not body.get("__pptx_content__"):
                return body
            if last_text and not _PPTX_RE.search(last_text) and not body.get("__pptx_content__"):
                return body

            topic = _topic(last_text) or body.get("__pptx_topic__") or "Презентация"
            num = _slide_count(last_text)

            await self._emit(__event_emitter__, f"📊 Генерирую презентацию: «{topic[:50]}»…")

            # Accept pre-generated slide content from router to avoid double LLM call
            pre_content = body.get("__pptx_content__")
            if pre_content and isinstance(pre_content, list):
                slides = pre_content[:num]
            else:
                slides = await _gen_slides(
                    self.valves.base_url, self.valves.api_key,
                    self.valves.model, topic, num,
                )

            await self._emit(__event_emitter__, "🔨 Собираю .pptx…")
            out_dir = Path(self.valves.output_dir)
            out_dir.mkdir(parents=True, exist_ok=True)
            filename = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
            file_path = str(out_dir / filename)
            loop = asyncio.get_event_loop()
            await loop.run_in_executor(None, _build_pptx, topic, slides, file_path)

            download_url = f"{self.valves.download_base_url}/{filename}"
            html = _html_preview(topic, slides)

            reply = (
                f"📊 **Презентация готова** — {len(slides) + 1} слайдов\n\n"
                f"[⬇️ Скачать {filename}]({download_url})\n\n"
                f"```html\n{html}\n```"
            )

            await self._emit(__event_emitter__, "✅ Готово", done=True)

            body["messages"] = [
                {"role": "system", "content": "Выведи следующее сообщение пользователю дословно."},
                {"role": "user", "content": reply},
            ]
            body["model"] = self.valves.model
            body["stream"] = False
            body["__pptx_handled__"] = True

        except Exception as e:
            log.error("PPTX error: %s", e)
            await self._emit(__event_emitter__, f"❌ Ошибка: {e}", done=True)

        return body

    async def outlet(
        self,
        body: dict,
        __event_emitter__: Optional[Callable[[dict], Awaitable[None]]] = None,
        user: Optional[dict] = None,
        __user__: Optional[dict] = None,
    ) -> dict:
        return body


# OpenWebUI Pipelines compatibility alias
Pipeline = Pipeline
